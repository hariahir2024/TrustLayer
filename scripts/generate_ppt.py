import sys
import os
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Slide Dimensions (16:9 widescreen)
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

# Premium Color Palette (Deep Tech Slate & Cyan/Indigo Accents)
COLOR_BG = RGBColor(3, 7, 18)           # #030712 - Deep Dark Black
COLOR_CARD_BG = RGBColor(11, 19, 41)     # #0b1329 - Dark Card Blue
COLOR_TEXT_PRIMARY = RGBColor(255, 255, 255) # White
COLOR_TEXT_SECONDARY = RGBColor(148, 163, 184) # Slate gray (#94a3b8)
COLOR_ACCENT_CYAN = RGBColor(6, 182, 212)    # Cyan (#06b6d4)
COLOR_ACCENT_INDIGO = RGBColor(99, 102, 241) # Indigo (#6366f1)
COLOR_ACCENT_GREEN = RGBColor(16, 185, 129)  # Green (#10b981)
COLOR_ACCENT_RED = RGBColor(239, 68, 68)     # Red (#ef4444)

def set_slide_background(slide):
    """Set slide background to solid dark color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text=""):
    """Create a unified top header bar for content slides."""
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.name = "Inter"
        cat_p.font.size = Pt(10)
        cat_p.font.bold = True
        cat_p.font.color.rgb = COLOR_ACCENT_CYAN
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.name = "Inter"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Add a thin accent line under header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT_INDIGO
    line.line.color.rgb = COLOR_ACCENT_INDIGO

def add_footer(slide, current_slide, total_slides):
    """Add a professional bottom footer with slide numbering."""
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
    footer_tf = footer_box.text_frame
    footer_tf.margin_left = footer_tf.margin_top = footer_tf.margin_right = footer_tf.margin_bottom = 0
    footer_p = footer_tf.paragraphs[0]
    footer_p.text = f"Team SOLARIS  •  CBI Hackathon 2026 Phase II  •  TrustLayer Prototype"
    footer_p.font.name = "Inter"
    footer_p.font.size = Pt(9)
    footer_p.font.color.rgb = COLOR_TEXT_SECONDARY
    
    num_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.3))
    num_tf = num_box.text_frame
    num_tf.margin_left = num_tf.margin_top = num_tf.margin_right = num_tf.margin_bottom = 0
    num_p = num_tf.paragraphs[0]
    num_p.text = f"{current_slide} / {total_slides}"
    num_p.font.name = "Inter"
    num_p.font.size = Pt(9)
    num_p.font.color.rgb = COLOR_TEXT_SECONDARY
    num_p.alignment = PP_ALIGN.RIGHT

def create_card_shape(slide, left, top, width, height, title="", border_color=None):
    """Add a subtle card shape with border and solid fill."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.color.rgb = RGBColor(30, 41, 59) # Slate slate border (#1e293b)
        card.line.width = Pt(1)
    return card

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_layout = prs.slide_layouts[6]
    
    total_slides = 12

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # SVG Center decorative accent (large thin ring)
    ring = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.66), Inches(1.75), Inches(4.0), Inches(4.0))
    ring.fill.background()
    ring.line.color.rgb = COLOR_ACCENT_CYAN
    ring.line.width = Pt(1)

    # Title box
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_p = t_tf.paragraphs[0]
    t_p.text = "TrustLayer"
    t_p.font.name = "Inter"
    t_p.font.size = Pt(64)
    t_p.font.bold = True
    t_p.font.color.rgb = COLOR_TEXT_PRIMARY
    t_p.alignment = PP_ALIGN.CENTER

    t_sub = t_tf.add_paragraph()
    t_sub.text = "Continuous AI Behavioral Biometric Authentication"
    t_sub.font.name = "Inter"
    t_sub.font.size = Pt(20)
    t_sub.font.color.rgb = COLOR_ACCENT_CYAN
    t_sub.alignment = PP_ALIGN.CENTER
    t_sub.space_before = Pt(15)

    # Bottom team detail
    det_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.0))
    det_tf = det_box.text_frame
    det_p = det_tf.paragraphs[0]
    det_p.text = "Team SOLARIS  •  MNNIT Allahabad"
    det_p.font.name = "Inter"
    det_p.font.size = Pt(15)
    det_p.font.bold = True
    det_p.font.color.rgb = COLOR_TEXT_PRIMARY
    det_p.alignment = PP_ALIGN.CENTER
    
    det_p2 = det_tf.add_paragraph()
    det_p2.text = "Cyber Security Hackathon 2026 Phase II  •  Bharat Suraksha Bank"
    det_p2.font.name = "Inter"
    det_p2.font.size = Pt(12)
    det_p2.font.color.rgb = COLOR_TEXT_SECONDARY
    det_p2.alignment = PP_ALIGN.CENTER
    det_p2.space_before = Pt(8)

    add_footer(slide1, 1, total_slides)

    # ==========================================
    # SLIDE 2: Team Introduction
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Meet Team SOLARIS", "Introduction")
    
    # Left Column: Team Cards (Vertical List)
    y_start = Inches(1.8)
    card_w = Inches(5.5)
    card_h = Inches(1.4)
    gap = Inches(0.2)
    
    team_members = [
        {"name": "Hari Ahir", "role": "Lead Backend & ML Engineer", "detail": "Implemented PyTorch LSTMs, XGBoost fusion, SQLite schema, & lockout backend."},
        {"name": "Team Solaris Member 2", "role": "Frontend & SDK Developer", "detail": "Created static/js/sdk.js telemetry recorder, bank.html interface & dashboard.js chart feeds."},
        {"name": "Team Solaris Member 3", "role": "Security Researcher & QA", "detail": "Analyzed CMU/Balabit telemetry distributions, KMT datasets, & wrote verify_integration.py."}
    ]

    for idx, mem in enumerate(team_members):
        y_pos = y_start + idx * (card_h + gap)
        create_card_shape(slide2, Inches(0.8), y_pos, card_w, card_h)
        
        # Profile icon block
        icon_box = slide2.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.15), Inches(0.8), Inches(0.8))
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = "👤"
        icon_p.font.size = Pt(28)
        
        # Text details
        txt_box = slide2.shapes.add_textbox(Inches(1.9), y_pos + Inches(0.15), Inches(4.2), Inches(1.1))
        txt_tf = txt_box.text_frame
        txt_tf.word_wrap = True
        txt_tf.margin_left = txt_tf.margin_top = txt_tf.margin_right = txt_tf.margin_bottom = 0
        
        p1 = txt_tf.paragraphs[0]
        p1.text = mem["name"]
        p1.font.name = "Inter"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        p2 = txt_tf.add_paragraph()
        p2.text = mem["role"]
        p2.font.name = "Inter"
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_ACCENT_CYAN
        p2.space_before = Pt(3)
        
        p3 = txt_tf.add_paragraph()
        p3.text = mem["detail"]
        p3.font.name = "Inter"
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = COLOR_TEXT_SECONDARY
        p3.space_before = Pt(4)

    # Right Column: Big quote
    q_box = slide2.shapes.add_textbox(Inches(6.8), Inches(2.6), Inches(5.7), Inches(3.0))
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    q_tf.margin_left = q_tf.margin_top = q_tf.margin_right = q_tf.margin_bottom = 0
    
    qp = q_tf.paragraphs[0]
    qp.text = "“"
    qp.font.name = "Georgia"
    qp.font.size = Pt(72)
    qp.font.bold = True
    qp.font.color.rgb = COLOR_ACCENT_CYAN
    
    qp2 = q_tf.add_paragraph()
    qp2.text = "We set out to answer one critical question:\nIf an attacker already has your NetBanking password — what stops them?"
    qp2.font.name = "Inter"
    qp2.font.size = Pt(22)
    qp2.font.bold = True
    qp2.font.color.rgb = COLOR_TEXT_PRIMARY
    qp2.space_before = Pt(10)
    
    qp3 = q_tf.add_paragraph()
    qp3.text = "Our solution is TrustLayer: continuous invisible behavioral validation, blocking account takeovers at the point of action."
    qp3.font.name = "Inter"
    qp3.font.size = Pt(14)
    qp3.font.color.rgb = COLOR_TEXT_SECONDARY
    qp3.space_before = Pt(15)

    add_footer(slide2, 2, total_slides)

    # ==========================================
    # SLIDE 3: Problem Statement
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "The Mid-Session Security Gap", "The Problem")

    # Left box: The timeline gap
    create_card_shape(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    t_box = slide3.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(4.9), Inches(4.2))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_tf.margin_left = t_tf.margin_top = t_tf.margin_right = t_tf.margin_bottom = 0
    
    p = t_tf.paragraphs[0]
    p.text = "Current Authentication Strategy:"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2 = t_tf.add_paragraph()
    p2.text = "🟢 LOGIN CHECK (Static & Weak)\nPassword + OTP verification occurs only once at the gate."
    p2.font.name = "Inter"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_ACCENT_GREEN
    p2.space_before = Pt(15)
    
    p3 = t_tf.add_paragraph()
    p3.text = "🔴 MID-SESSION BLIND SPOT (The Gap)\nAfter login, the session is trusted unconditionally. If a session token is hijacked, or a device is left open, the bank is blind."
    p3.font.name = "Inter"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_ACCENT_RED
    p3.space_before = Pt(15)
    
    p4 = t_tf.add_paragraph()
    p4.text = "🚨 ATTACK VECTORS CAPTURED:\n• Account Takeover (Session hijacking)\n• Credential Sharing (Unauthorized logins)\n• Automated Bots (Programmatic credential stuffing)"
    p4.font.name = "Inter"
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_TEXT_SECONDARY
    p4.space_before = Pt(20)

    # Right: Stat grid
    stat_data = [
        {"stat": "₹29,082 Cr", "label": "Banking Fraud FY23-24", "detail": "Reported by RBI Annual Report"},
        {"stat": "47% YoY", "label": "Growth in Digital Frauds", "detail": "Rising threat vector in digital channels"},
        {"stat": "34% attacks", "label": "Via Automated Scripts", "detail": "Bots and headless browsers bypassing front gate"}
    ]
    
    for idx, stat in enumerate(stat_data):
        y_pos = Inches(1.8) + idx * Inches(1.6)
        create_card_shape(slide3, Inches(6.8), y_pos, Inches(5.7), Inches(1.4))
        
        s_box = slide3.shapes.add_textbox(Inches(7.1), y_pos + Inches(0.15), Inches(5.1), Inches(1.1))
        s_tf = s_box.text_frame
        s_tf.word_wrap = True
        s_tf.margin_left = s_tf.margin_top = s_tf.margin_right = s_tf.margin_bottom = 0
        
        sp1 = s_tf.paragraphs[0]
        sp1.text = stat["stat"]
        sp1.font.name = "Inter"
        sp1.font.size = Pt(28)
        sp1.font.bold = True
        sp1.font.color.rgb = COLOR_ACCENT_CYAN
        
        sp2 = s_tf.add_paragraph()
        sp2.text = stat["label"]
        sp2.font.name = "Inter"
        sp2.font.size = Pt(12)
        sp2.font.bold = True
        sp2.font.color.rgb = COLOR_TEXT_PRIMARY
        sp2.space_before = Pt(4)
        
        sp3 = s_tf.add_paragraph()
        sp3.text = stat["detail"]
        sp3.font.name = "Inter"
        sp3.font.size = Pt(10)
        sp3.font.color.rgb = COLOR_TEXT_SECONDARY
        sp3.space_before = Pt(2)

    add_footer(slide3, 3, total_slides)

    # ==========================================
    # SLIDE 4: Solution Overview
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "TrustLayer — Continuous Behavioral Trust", "Solution")

    # Center diagram: Flow
    create_card_shape(slide4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.8))
    
    flow_box = slide4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(1.4))
    flow_tf = flow_box.text_frame
    flow_tf.margin_left = flow_tf.margin_top = flow_tf.margin_right = flow_tf.margin_bottom = 0
    
    fp = flow_tf.paragraphs[0]
    fp.text = "How TrustLayer Secures the Session:"
    fp.font.name = "Inter"
    fp.font.size = Pt(16)
    fp.font.bold = True
    fp.font.color.rgb = COLOR_TEXT_PRIMARY
    
    fp2 = flow_tf.add_paragraph()
    fp2.text = "⌨️ USER TYPING (Keystroke timings)   ──▶   🖱️ MOUSE NAVIGATION (Trajectories)   ──▶   🧠 3-MODEL ML FUSION   ──▶   🚦 DYNAMIC BLOCK/CHALLENGE"
    fp2.font.name = "Inter"
    fp2.font.size = Pt(12.5)
    fp2.font.bold = True
    fp2.font.color.rgb = COLOR_ACCENT_CYAN
    fp2.space_before = Pt(15)
    
    fp3 = flow_tf.add_paragraph()
    fp3.text = "Runs silently in the browser. Zero friction for legitimate users. No tokens or SMS required unless anomaly detected."
    fp3.font.name = "Inter"
    fp3.font.size = Pt(11)
    fp3.font.color.rgb = COLOR_TEXT_SECONDARY
    fp3.space_before = Pt(10)

    # Lower Pillars
    pillars = [
        {"title": "👁️ Invisible Security", "desc": "Passive telemetry monitors keystroke flight times, mouse velocity, and screen interaction curves without interrupting the banking flow."},
        {"title": "🔒 Privacy By Design", "desc": "Raw keystrokes or coordinates never leave the user's browser. The SDK processes events locally and transmits only normalized timings (μ/σ) to the backend."},
        {"title": "🔌 Drop-In SDK Integration", "desc": "Can be embedded into any digital portal by adding a single JavaScript SDK line. No infrastructure refactoring or hardware token requirements."}
    ]
    
    for idx, pil in enumerate(pillars):
        x_pos = Inches(0.8) + idx * Inches(4.0)
        create_card_shape(slide4, x_pos, Inches(3.9), Inches(3.7), Inches(2.7))
        
        p_box = slide4.shapes.add_textbox(x_pos + Inches(0.2), Inches(4.1), Inches(3.3), Inches(2.3))
        p_tf = p_box.text_frame
        p_tf.word_wrap = True
        p_tf.margin_left = p_tf.margin_top = p_tf.margin_right = p_tf.margin_bottom = 0
        
        pp1 = p_tf.paragraphs[0]
        pp1.text = pil["title"]
        pp1.font.name = "Inter"
        pp1.font.size = Pt(16)
        pp1.font.bold = True
        pp1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        pp2 = p_tf.add_paragraph()
        pp2.text = pil["desc"]
        pp2.font.name = "Inter"
        pp2.font.size = Pt(11.5)
        pp2.font.color.rgb = COLOR_TEXT_SECONDARY
        pp2.space_before = Pt(10)

    add_footer(slide4, 4, total_slides)

    # ==========================================
    # SLIDE 5: System Architecture
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "System Architecture Map", "Architecture")

    # Draw client backend dashboard layout
    create_card_shape(slide5, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=COLOR_ACCENT_CYAN)
    create_card_shape(slide5, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=COLOR_ACCENT_INDIGO)
    create_card_shape(slide5, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=COLOR_ACCENT_GREEN)

    # Card 1: Browser SDK
    c1 = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.3), Inches(4.4))
    c1_tf = c1.text_frame
    c1_tf.word_wrap = True
    c1_tf.margin_left = c1_tf.margin_top = c1_tf.margin_right = c1_tf.margin_bottom = 0
    p = c1_tf.paragraphs[0]
    p.text = "CLIENT LAYER (SDK)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_CYAN
    p2 = c1_tf.add_paragraph()
    p2.text = "📁 static/js/sdk.js\n📁 static/js/bank.js\n\n• Listens to keydown/keyup timestamps on username + password.\n• Records mouse coordinates dynamically.\n• Extracts browser metadata (screen resolution, timezone, language).\n• Signs requests with incremental nonces & timestamps."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)

    # Card 2: Backend API
    c2 = slide5.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.3), Inches(4.4))
    c2_tf = c2.text_frame
    c2_tf.word_wrap = True
    c2_tf.margin_left = c2_tf.margin_top = c2_tf.margin_right = c2_tf.margin_bottom = 0
    p = c2_tf.paragraphs[0]
    p.text = "BACKEND API LAYER"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_INDIGO
    p2 = c2_tf.add_paragraph()
    p2.text = "📁 app.py (FastAPI)\n📁 db_sqlite.py (SQLite)\n\n• Exposes REST endpoints for auth, score, registration & re-auth.\n• Validates CORS, IP blocks, and token authenticity.\n• Implements persistent lockout tracking via database.\n• SQLite runs in WAL mode for concurrently handling score updates."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)

    # Card 3: ML Engine
    c3 = slide5.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.4))
    c3_tf = c3.text_frame
    c3_tf.word_wrap = True
    c3_tf.margin_left = c3_tf.margin_top = c3_tf.margin_right = c3_tf.margin_bottom = 0
    p = c3_tf.paragraphs[0]
    p.text = "ML SCORING ENGINE"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_GREEN
    p2 = c3_tf.add_paragraph()
    p2.text = "📁 ml_engine.py (AI Engine)\n📁 models/ (Pre-trained files)\n\n• Evaluates keystroke reconstruction error via PyTorch LSTM.\n• Evaluates mouse trajectories anomaly via PyTorch LSTM.\n• Combines timings, device metrics, & time-of-day history.\n• Fuses scores using XGBoost Classifier with SHAP explainability."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)

    add_footer(slide5, 5, total_slides)

    # ==========================================
    # SLIDE 6: ML Pipeline & Datasets
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "The Three-Model ML Pipeline", "ML Engine")

    # Left: Vertical flow cards
    flow_steps = [
        {"title": "🤖 Heuristic pre-screen", "desc": "Instant bot filters: checks WebDriver flag, zero velocity cursor, 0ms flight time. Bypasses ML on bot detections to save CPU."},
        {"title": "🧠 LSTM Autoencoders", "desc": "Double neural networks (Keystroke & Mouse). Trained on Balabit & CMU datasets to reconstruct sequence timings. High MSE = anomaly."},
        {"title": "🌲 XGBoost Fusion Classifier", "desc": "Blends Keystroke score, Mouse score, Device Class risk, and Time-of-day history. Predicts fraud probabilities with SHAP features."}
    ]
    for idx, step in enumerate(flow_steps):
        y_pos = Inches(1.8) + idx * Inches(1.6)
        create_card_shape(slide6, Inches(0.8), y_pos, Inches(5.5), Inches(1.4))
        
        s_box = slide6.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.15), Inches(5.1), Inches(1.1))
        s_tf = s_box.text_frame
        s_tf.word_wrap = True
        s_tf.margin_left = s_tf.margin_top = s_tf.margin_right = s_tf.margin_bottom = 0
        
        sp1 = s_tf.paragraphs[0]
        sp1.text = step["title"]
        sp1.font.name = "Inter"
        sp1.font.size = Pt(14)
        sp1.font.bold = True
        sp1.font.color.rgb = COLOR_ACCENT_CYAN
        
        sp2 = s_tf.add_paragraph()
        sp2.text = step["desc"]
        sp2.font.name = "Inter"
        sp2.font.size = Pt(10)
        sp2.font.color.rgb = COLOR_TEXT_SECONDARY
        sp2.space_before = Pt(6)

    # Right: Dataset Performance
    create_card_shape(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    
    perf_box = slide6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
    perf_tf = perf_box.text_frame
    perf_tf.word_wrap = True
    perf_tf.margin_left = perf_tf.margin_top = perf_tf.margin_right = perf_tf.margin_bottom = 0
    
    p = perf_tf.paragraphs[0]
    p.text = "Dataset & Training Metrics:"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2 = perf_tf.add_paragraph()
    p2.text = "📊 Trained on KMT Dataset (88 users, real banking sessions)\n• XGBoost Fusion Accuracy: 86.87%\n• Intruder Detection Recall: 92.40% (high safety index)\n• Precision: 79.41%  • F1-Score: 85.41%\n• Execution Latency: < 2ms (ideal for real-time gating)"
    p2.font.name = "Inter"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)
    
    p3 = perf_tf.add_paragraph()
    p3.text = "🎯 Feature Importance Distribution (SHAP):"
    p3.font.name = "Inter"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_TEXT_PRIMARY
    p3.space_before = Pt(20)
    
    p4 = perf_tf.add_paragraph()
    p4.text = "• Keystroke Cadence: 45.36% contribution\n• Metadata & Time-of-day check: 24.98% contribution\n• Mouse Cursor Velocity: 20.85% contribution\n• Device Fingerprint Match: 8.81% contribution"
    p4.font.name = "Inter"
    p4.font.size = Pt(11)
    p4.font.color.rgb = COLOR_ACCENT_CYAN
    p4.space_before = Pt(10)

    add_footer(slide6, 6, total_slides)

    # ==========================================
    # SLIDE 7: 7-Band Risk Escalation System
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Proportional Security — The 7-Band Scale", "Risk Escalation")

    # Draw horizontal scale
    bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_CARD_BG
    bar.line.color.rgb = RGBColor(30, 41, 59)

    # Accent sections inside the bar
    sect1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.5), Inches(0.4))
    sect1.fill.solid()
    sect1.fill.fore_color.rgb = COLOR_ACCENT_GREEN
    sect1.line.fill.background()
    
    sect2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), Inches(1.8), Inches(4.7), Inches(0.4))
    sect2.fill.solid()
    sect2.fill.fore_color.rgb = COLOR_ACCENT_CYAN
    sect2.line.fill.background()

    sect3 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(1.8), Inches(3.5), Inches(0.4))
    sect3.fill.solid()
    sect3.fill.fore_color.rgb = COLOR_ACCENT_RED
    sect3.line.fill.background()

    # Labels below the scale bar
    labels_box = slide7.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.5))
    labels_tf = labels_box.text_frame
    labels_tf.margin_left = labels_tf.margin_top = labels_tf.margin_right = labels_tf.margin_bottom = 0
    lp = labels_tf.paragraphs[0]
    lp.text = "0  ─────────  30 (GREEN)  ─────────  45 (AMBER_L)  ─────  70 (AMBER_H)  ──────  82 (RED_H)  ─────  100 (RED_C)"
    lp.font.name = "Inter"
    lp.font.bold = True
    lp.font.size = Pt(11)
    lp.font.color.rgb = COLOR_TEXT_PRIMARY
    lp.alignment = PP_ALIGN.CENTER

    # Action pillars
    bands_desc = [
        {"title": "🟢 GREEN (0-30)", "desc": "Active session continues silently. Telemetry interval set to standard 15 seconds. Minimal scoring overhead."},
        {"title": "🟡 AMBER (31-70)", "desc": "Soft step-up re-auth modal (types baseline). Limits high-value transactions. Polling drops to 10 seconds."},
        {"title": "🔴 RED (71-100)", "desc": "Instant UI overlay freeze. Session invalidated. SMS/ntfy pushes sent. Telemetry interval accelerated to 5 seconds."}
    ]
    for idx, b_desc in enumerate(bands_desc):
        x_pos = Inches(0.8) + idx * Inches(4.0)
        create_card_shape(slide7, x_pos, Inches(3.2), Inches(3.7), Inches(3.4))
        
        b_box = slide7.shapes.add_textbox(x_pos + Inches(0.2), Inches(3.4), Inches(3.3), Inches(3.0))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = b_tf.margin_top = b_tf.margin_right = b_tf.margin_bottom = 0
        
        bp1 = b_tf.paragraphs[0]
        bp1.text = b_desc["title"]
        bp1.font.name = "Inter"
        bp1.font.size = Pt(16)
        bp1.font.bold = True
        bp1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        bp2 = b_tf.add_paragraph()
        bp2.text = b_desc["desc"]
        bp2.font.name = "Inter"
        bp2.font.size = Pt(11.5)
        bp2.font.color.rgb = COLOR_TEXT_SECONDARY
        bp2.space_before = Pt(12)

    add_footer(slide7, 7, total_slides)

    # ==========================================
    # SLIDE 8: Key Features & Innovation
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "What Makes TrustLayer Different", "Features")

    features_list = [
        {"title": "🔍 Explainable AI", "desc": "SHAP-lite feature attribution converts weights to plain English messages: 'Key flight time G→H deviated 3.4σ from enrolled mean'."},
        {"title": "🛡️ Bot Pre-Screening", "desc": "Heuristic filters intercept selenium/playwright automation immediately, flagging bot signatures before expensive ML pipelines run."},
        {"title": "❄️ Cold Start Solver", "desc": "Blends global CMU/Balabit population priors with user-specific data using Exponential Moving Average, maturing at 15 sessions."},
        {"title": "🔒 No Biometric Storage", "desc": "Complies with RBI storage guidelines. Processes key timings in client browser. Stores only statistical features (μ/σ)."},
        {"title": "🔁 Replay Protection", "desc": "Telemetry API payloads signed with UUID nonces and 60-second timestamp windows to prevent capture-and-replay attacks."},
        {"title": "📊 SOC Incident Feed", "desc": "WebSockets push risk alerts and feature breakdowns to a live security console. Analysts can force-freeze or dismiss sessions."}
    ]

    for idx, feat in enumerate(features_list):
        row = idx // 3
        col = idx % 3
        
        x_pos = Inches(0.8) + col * Inches(4.0)
        y_pos = Inches(1.8) + row * Inches(2.5)
        
        create_card_shape(slide8, x_pos, y_pos, Inches(3.7), Inches(2.2))
        
        f_box = slide8.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.15), Inches(3.3), Inches(1.9))
        f_tf = f_box.text_frame
        f_tf.word_wrap = True
        f_tf.margin_left = f_tf.margin_top = f_tf.margin_right = f_tf.margin_bottom = 0
        
        fp1 = f_tf.paragraphs[0]
        fp1.text = feat["title"]
        fp1.font.name = "Inter"
        fp1.font.size = Pt(14.5)
        fp1.font.bold = True
        fp1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        fp2 = f_tf.add_paragraph()
        fp2.text = feat["desc"]
        fp2.font.name = "Inter"
        fp2.font.size = Pt(10.5)
        fp2.font.color.rgb = COLOR_TEXT_SECONDARY
        fp2.space_before = Pt(8)

    add_footer(slide8, 8, total_slides)

    # ==========================================
    # SLIDE 9: Mockups & Visuals Placeholder
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Live Prototype Visual Interface", "Interactive Demo")

    # Left box: Bank portal mockup placeholder
    create_card_shape(slide9, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    b_mock = slide9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(4.9), Inches(4.4))
    b_tf = b_mock.text_frame
    b_tf.word_wrap = True
    b_tf.margin_left = b_tf.margin_top = b_tf.margin_right = b_tf.margin_bottom = 0
    p = b_tf.paragraphs[0]
    p.text = "🏛️ BHARAT SURAKSHA BANK PORTAL"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_CYAN
    p2 = b_tf.add_paragraph()
    p2.text = "A fully interactive internet banking mockup:\n\n• Dynamic Passphrase Generation: custom 11-char string based on first & last name suffix to maintain LSTM shape.\n• Fund Transfer & UPI Panel: gates transactions through risk rules.\n• Floating Demo Controller: lets examiners trigger 3 attack personas instantly:\n  1. Legitimate user (GREEN - continue)\n  2. Human Intruder (RED_HIGH - frozen overlay)\n  3. Automated Bot (RED_CRITICAL - locked)\n• Real-Time push notifications sent to user's phone via ntfy.sh topic."
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)

    # Right box: SOC Dashboard mockup placeholder
    create_card_shape(slide9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    soc_mock = slide9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
    soc_tf = soc_mock.text_frame
    soc_tf.word_wrap = True
    soc_tf.margin_left = soc_tf.margin_top = soc_tf.margin_right = soc_tf.margin_bottom = 0
    p = soc_tf.paragraphs[0]
    p.text = "📊 FRAUD OPS SECURITY DASHBOARD"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_INDIGO
    p2 = soc_tf.add_paragraph()
    p2.text = "The live security console for banking fraud teams:\n\n• WebSocket Real-Time feed: score changes push updates to timeline without page refresh.\n• Feature Attributions Panel: visualizes exactly which keystroke positions or mouse intervals caused risk peaks.\n• Analyst Overrides: 'Force Freeze' to lock compromised session, or 'Mark False Positive' to restore access & downgrade threat score.\n• Frozen Sessions Tab: shows full details of all auto-frozen, manually blocked, and bot-blocked sessions."
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)

    add_footer(slide9, 9, total_slides)

    # ==========================================
    # SLIDE 10: Security Architecture & Compliance
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Enterprise-Grade Security Standards", "Security Center")

    # Left: Security features list
    create_card_shape(slide10, Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.8))
    s_box = slide10.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.6), Inches(4.4))
    s_tf = s_box.text_frame
    s_tf.word_wrap = True
    s_tf.margin_left = s_tf.margin_top = s_tf.margin_right = s_tf.margin_bottom = 0
    
    p = s_tf.paragraphs[0]
    p.text = "Core Protection Layers:"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2 = s_tf.add_paragraph()
    p2.text = "• Credential Hashing: PBKDF2-SHA256 with 100k iterations and randomized salt. Wipes out dictionary attack vectors.\n\n• Persistent Rate Limiter: Blocks IP addresses sending >5 incorrect passwords within 60 seconds (prevents brute-force).\n\n• Robust Session Lockout: SQLite-backed table persists consecutive failures (password or biometrics) across server reboots. Lock duration is 15 minutes.\n\n• Dynamic Transaction Gating: Soft limits (₹5,000) trigger step-up verification. Hard limits (₹10,000) block immediately on compromised sessions."
    p2.font.name = "Inter"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.space_before = Pt(15)

    # Right: Standards compliance card grid
    compliance_items = [
        {"title": "NIST SP 800-63B", "desc": "Authentication Guidelines"},
        {"title": "OWASP API Top 10", "desc": "Injection & Rate-limit guards"},
        {"title": "RBI Digital Payments", "desc": "Security directive compliant"},
        {"title": "IT Act 2000 (Data Privacy)", "desc": "Stat-only storage design"}
    ]
    for idx, comp in enumerate(compliance_items):
        row = idx // 2
        col = idx % 2
        x_pos = Inches(7.3) + col * Inches(2.7)
        y_pos = Inches(1.8) + row * Inches(2.3)
        
        create_card_shape(slide10, x_pos, y_pos, Inches(2.5), Inches(2.1))
        
        c_box = slide10.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.15), Inches(2.2), Inches(1.8))
        c_tf = c_box.text_frame
        c_tf.word_wrap = True
        c_tf.margin_left = c_tf.margin_top = c_tf.margin_right = c_tf.margin_bottom = 0
        
        cp1 = c_tf.paragraphs[0]
        cp1.text = comp["title"]
        cp1.font.name = "Inter"
        cp1.font.size = Pt(14)
        cp1.font.bold = True
        cp1.font.color.rgb = COLOR_ACCENT_CYAN
        
        cp2 = c_tf.add_paragraph()
        cp2.text = comp["desc"]
        cp2.font.name = "Inter"
        cp2.font.size = Pt(10.5)
        cp2.font.color.rgb = COLOR_TEXT_SECONDARY
        cp2.space_before = Pt(10)

    add_footer(slide10, 10, total_slides)

    # ==========================================
    # SLIDE 11: Challenges Encountered
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Challenges & How We Solved Them", "Challenges")

    challenges = [
        {"title": "⚡ COLD START CHALLENGE", "subtitle": "Biometric Enrollment Gap", "desc": "Enrolled profiles require baseline data. A first-time user has no history, making authentic authentication impossible on session 1.\n\n✔ Solution: Hybrid Prior Blending. Blended population mean (CMU/Balabit benchmarks) with incoming user samples using Exponential Moving Average. Full transition completed at session 15."},
        {"title": "⚡ DATASET BIAS CHALLENGE", "subtitle": "Western Keyboard Prior Bias", "desc": "Academic biometrics datasets are collected on Western populations with specific layout behaviors. Indian typing cadences (bilingual, different shorthands) vary.\n\n✔ Solution: KMT custom banking dataset consisting of 88 users with real billing input timings. XGBoost fusion was trained directly on this custom set, improving baseline accuracy to 86.87%."},
        {"title": "⚡ LOCALHOST BANS CHALLENGE", "subtitle": "Developer Lockout During Demos", "desc": "Auto-blocking IPs of bot threats works in production, but on localhost (127.0.0.1) it blocks the developer's own browser, breaking live presentations.\n\n✔ Solution: Loopback Bypass. Implemented a network guard in db_sqlite.py to detect loopback/localhost IPs. The system logs the threat event but bypasses the IP-block execution, keeping the demo running."}
    ]

    for idx, chal in enumerate(challenges):
        x_pos = Inches(0.8) + idx * Inches(4.0)
        create_card_shape(slide11, x_pos, Inches(1.8), Inches(3.7), Inches(4.8))
        
        ch_box = slide11.shapes.add_textbox(x_pos + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
        ch_tf = ch_box.text_frame
        ch_tf.word_wrap = True
        ch_tf.margin_left = ch_tf.margin_top = ch_tf.margin_right = ch_tf.margin_bottom = 0
        
        chp1 = ch_tf.paragraphs[0]
        chp1.text = chal["title"]
        chp1.font.name = "Inter"
        chp1.font.size = Pt(14)
        chp1.font.bold = True
        chp1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        chp2 = ch_tf.add_paragraph()
        chp2.text = chal["subtitle"]
        chp2.font.name = "Inter"
        chp2.font.size = Pt(11)
        chp2.font.bold = True
        chp2.font.color.rgb = COLOR_ACCENT_CYAN
        chp2.space_before = Pt(4)
        
        chp3 = ch_tf.add_paragraph()
        chp3.text = chal["desc"]
        chp3.font.name = "Inter"
        chp3.font.size = Pt(10.5)
        chp3.font.color.rgb = COLOR_TEXT_SECONDARY
        chp3.space_before = Pt(15)

    add_footer(slide11, 11, total_slides)

    # ==========================================
    # SLIDE 12: Impact & Roadmap
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "Impact & Future Development Roadmap", "Roadmap")

    # Top impact cards
    impacts = [
        {"stat": "8 Billion", "label": "Monthly UPI Transactions", "detail": "Every digital payment can be passively verified & protected."},
        {"stat": "₹29,082 Crore", "label": "Addressable Annual Fraud", "detail": "Mitigates credential theft losses across Indian banking networks."},
        {"stat": "92.4% Recall", "label": "High intruder block rate", "desc": "Catches 9 in 10 hijackers with zero friction for legit owners."}
    ]
    for idx, imp in enumerate(impacts):
        x_pos = Inches(0.8) + idx * Inches(4.0)
        create_card_shape(slide12, x_pos, Inches(1.8), Inches(3.7), Inches(1.8))
        
        i_box = slide12.shapes.add_textbox(x_pos + Inches(0.2), Inches(1.95), Inches(3.3), Inches(1.5))
        i_tf = i_box.text_frame
        i_tf.word_wrap = True
        i_tf.margin_left = i_tf.margin_top = i_tf.margin_right = i_tf.margin_bottom = 0
        
        ip1 = i_tf.paragraphs[0]
        ip1.text = imp["stat"]
        ip1.font.name = "Inter"
        ip1.font.size = Pt(22)
        ip1.font.bold = True
        ip1.font.color.rgb = COLOR_ACCENT_CYAN
        
        ip2 = i_tf.add_paragraph()
        ip2.text = imp["label"]
        ip2.font.name = "Inter"
        ip2.font.size = Pt(11)
        ip2.font.bold = True
        ip2.font.color.rgb = COLOR_TEXT_PRIMARY
        ip2.space_before = Pt(3)
        
        ip3 = i_tf.add_paragraph()
        ip3.text = imp.get("detail", imp.get("desc", ""))
        ip3.font.name = "Inter"
        ip3.font.size = Pt(9.5)
        ip3.font.color.rgb = COLOR_TEXT_SECONDARY
        ip3.space_before = Pt(2)

    # Bottom timeline roadmap
    create_card_shape(slide12, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.7))
    
    t_box = slide12.shapes.add_textbox(Inches(1.1), Inches(4.1), Inches(11.1), Inches(2.3))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_tf.margin_left = t_tf.margin_top = t_tf.margin_right = t_tf.margin_bottom = 0
    
    tp1 = t_tf.paragraphs[0]
    tp1.text = "Deployment Roadmap Phases:"
    tp1.font.name = "Inter"
    tp1.font.size = Pt(16)
    tp1.font.bold = True
    tp1.font.color.rgb = COLOR_TEXT_PRIMARY
    
    tp2 = t_tf.add_paragraph()
    tp2.text = "✔ [Phase 1 - NetBanking] NetBanking Web SDK with dynamic passphrase calibration & SQLite persistent lockout. (Done)\n" + \
               "⏳ [Phase 2 - Mobile SDK] Gyroscope, Accelerometer, and Touch Pressure behavioral capture on Native Android & iOS.\n" + \
               "⏳ [Phase 3 - UPI Layer] Passively analyzing Virtual Payment Address (VPA) entry rhythms and transaction speed velocity.\n" + \
               "⏳ [Phase 4 - Federated Learning] Passive opt-in Indian Population Prior model updates via secure on-device training."
    tp2.font.name = "Inter"
    tp2.font.size = Pt(11.5)
    tp2.font.color.rgb = COLOR_TEXT_SECONDARY
    tp2.space_before = Pt(10)

    add_footer(slide12, 12, total_slides)

    # Save presentation
    output_filename = "C:\\hackathon\\cbi hackathon\\TrustLayer_Presentation_Deck.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully to {output_filename}")

if __name__ == "__main__":
    main()
